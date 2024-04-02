from pptx import Presentation


from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt, Inches
from pptx.enum.text import MSO_AUTO_SIZE

from src.llm.lang import H2OGPTE, query, try_and_parse, Session

def gen_key_words(client: H2OGPTE, user_prompt: str, llm='mistralai/Mixtral-8x7B-Instruct-v0.1') -> list[str]:
    sys = """You are an assistant whose task is to perform searches on the internet on a specific topic.\
    The user is interested to create a presentation about a topic of interest.\
    Think about what to do, then reply with your thought process and at least one corresponding google query as an array in JSON format,\
    but limit yourself to 5 queries.\
    The JSON array should be contained in a code chunk. Keep strictly to the format in the example below.
    Example:
    Since the presentation is about milk, I will probably want to search up the different types of milk, ...
    ```json
    ["Milk", "Oat Milk", "Plant-based milks", "Cow Milk", "Goat Milk"]
    ```
    """
    user = user_prompt

    
    return try_and_parse(client, user, sys, llm=llm)

def format_into_readable_list(ls: list[str]) -> str:
    for i in range(len(ls)):
        ls[i] = f"{i+1}." + ls[i] 
    return "\n".join(ls)

def decide_slide_titles(client: H2OGPTE, presentation_desc: str, information: list[str]) -> list[str]:
    sys = f"""You are an assistant whose task is to help a user in creating a presentation. {presentation_desc}. 
    Here are some papers/article titles to help you decide on slide titles:\n{format_into_readable_list(information)}"""

    user = """You are provided with a list of article/paper titles. 
    Using these articles, you will create a presentation. 
    1. Explain how you would design the presentation slides such that the presentation will flow well.
    2. Think of a suitable title for this presentation.
    3. Then, plan your slides and give me your slide titles in a json array within a code chunk. Include the title slide, which is the title for the presentation.
    Adhere strictly to the following example:
    1. I would introduce the Transformers franchise and provide general information about its history to ease my viewers into the subject. \
        Then, I will think about subtopics, such as the Transformers films, Transformers characters and Transformers in comics, using the wikipedia entry summaries\
        I have been provided. 
    2. I think a good title for this presentation is "Transformers: An Overview".
    3. Here is the json array of slide titles:
    ```json
        [
        "Transformers: An Overview", 
        "Introduction to Transformers", 
        "Transformers in Film",
        "Transformers: Revenge of the Fallen (2009)", 
        "Bumblebee (2018)", 
        "Characters in the Transformers Universe", 
        "Transformers in comics",
        "Conclusion"
        ]
    ```
    """
    return try_and_parse(client, user, sys)


def decide_ppt_colour(client: H2OGPTE, presentation_desc: str) -> list[dict]:
    formatted = f"""{presentation_desc} Think of a good background colour, in RGB format,\
        for the slides and a good colour, also in RGB format, for the\
        text. Typically, if the text colour is bright (for example RGB [255, 255, 255] is white), then the background colour should be dark
        (RGB [0, 0, 100] is dark blue). Conversely, if the text colour is dark (for example RGB [0, 0, 0] is black), the background colour should be bright\
        . You are free to choose any text and background colour, \
        as long as you follow these rules. Please do not assign grey-scale colours for the text and background (like RGB [50, 50, 50]), as much as possible.

        Explain clearly why you chose the background and text colours. Then, generate a code chunk. Within the code chunk,\
        provide a JSON array containing two colours. Do not say anything else. Adhere strictly to the example reply below:
        I chose blue RGB [0, 35, 140] for the background color and light yellow RGB [255, 234, 0] for the font color. The contrast makes it easy to read.\
        Furthermore, the colours blue and yellow are associated with the Pokémon Franchise.
        ```
        [{{"background": [0, 0, 140]}}, {{"text": [255, 234, 0]}}]
        ```
    """

    sys = "You are an assistant whose task is to help a user in creating a presentation."
    return try_and_parse(client, formatted, system_prompt=sys)


def slide_query(session: Session, slide: str, llm="mistralai/Mixtral-8x7B-Instruct-v0.1"):
    output = session.query(
        rag_config={
                "rag_type": "hyde1",
        },
        llm=llm,
        system_prompt="You are an assistant whose task is to help a user in creating a presentation. ",
        pre_prompt_query="""You have been provided with the following information, which may be useful in your task. 
        Whenever the user gives you a task, summarise your findings into short sentences and keep your reply to 2 paragraphs""",
        prompt_query="""Decide if the information is relevant, and use it if needed""",
        message=f"""Generate some content about {slide}. """,
        
    )
    return output.content

def generate_ppt(
        client: H2OGPTE, 
        chat_session_id, 
        slide_titles: list[str],
        presentation_desc: str, 
        user_decided_colours: dict=dict(),
        llm: str="mistralai/Mixtral-8x7B-Instruct-v0.1"
        ):
    
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    if not user_decided_colours:
        colours = decide_ppt_colour(client, presentation_desc)
        background = RGBColor(
            *tuple(list(
                colours[0].values()        
            )[0])
        )
        
        font = RGBColor(*tuple(list(colours[1].values())[0]))

    fill = title_slide.background.fill
    fill.solid()
    fill.fore_color.rgb = background
    title_slide.shapes.title.text = slide_titles[0]
    title_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb =  font
    title_slide.shapes.title.text_frame.paragraphs[0].font.name = 'Montserrat'
    title_slide.shapes.title.text_frame.paragraphs[0].font.bold = True


    first_shape =  title_slide.shapes[0]
    first_shape.left, first_shape.top, first_shape.width, first_shape.height = (prs.slide_width - Inches(12))//2, \
    (prs.slide_height-first_shape.height)//2 - Inches(1), Inches(12), Inches(2)



    with client.connect(chat_session_id) as session:
        for section in slide_titles:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = background

         
            contents = slide.placeholders[1]
            contents.text_frame.word_wrap = True

            title = slide.shapes.title
            title.text = section
            title.text_frame.paragraphs[0].font.color.rgb = font
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.name = 'Karla'
            # content != contents
            content = slide_query(session, section, llm=llm)
            contents.text = content

            for paragraph in contents.text_frame.paragraphs:
                paragraph.space_after = 1
                paragraph.space_before = 1
                # paragraph.level = 0
            
                paragraph.font.size = Pt(22)  
                paragraph.font.color.rgb = font
                paragraph.font.name = 'Karla'

            contents.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            shapes = slide.shapes
            new_width = Inches(14)
            new_height = Inches(7)
            shapes[0].height, shapes[0].width, shapes[0].top, shapes[0].left = shapes[0].height, new_width, shapes[0].top, (prs.slide_width-new_width)//2
            shapes[1].height, shapes[1].width, shapes[1].top, shapes[1].left = new_height, new_width, shapes[1].top, (prs.slide_width-new_width)//2
    
    return prs
        
        






